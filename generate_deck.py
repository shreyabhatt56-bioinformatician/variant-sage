import collections
import collections.abc
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_deck():
    prs = Presentation()
    
    # Set slide size to Widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Colors
    NAVY = RGBColor(15, 23, 42)
    BLUE = RGBColor(37, 99, 235)
    TEAL = RGBColor(13, 148, 136)
    WHITE = RGBColor(255, 255, 255)
    GRAY = RGBColor(100, 116, 139)
    LIGHT_BG = RGBColor(241, 245, 249)

    # --- 1. TITLE SLIDE ---
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Background
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    # Title text
    txbox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(2))
    tf = txbox.text_frame
    p = tf.add_paragraph()
    p.text = "Variant Sage"
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    p2 = tf.add_paragraph()
    p2.text = "AI-Powered Clinical Genomics Platform"
    p2.font.size = Pt(36)
    p2.font.color.rgb = TEAL

    p3 = tf.add_paragraph()
    p3.text = "Prepared by MolSys Team"
    p3.font.size = Pt(20)
    p3.font.color.rgb = GRAY
    p3.space_before = Pt(40)

    # --- 2. THE PROBLEM ---
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BG
    bg.line.fill.background()

    txbox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.3), Inches(1))
    tf = txbox.text_frame
    p = tf.add_paragraph()
    p.text = "The Bottleneck in Modern Genomics"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Three problem boxes
    probs = [
        ("Data Avalanche", "Next-Gen Sequencing (NGS) produces massive VCF files with millions of variants. Filtering them is a manual, error-prone nightmare."),
        ("Counselor Shortage", "Highly trained Genetic Counselors spend hours on literature review for a single patient. There are simply not enough experts."),
        ("Slow Turnaround", "Patients wait weeks for critical oncology or rare-disease reports. In clinical settings, time is literally life.")
    ]
    
    for i, (title, desc) in enumerate(probs):
        left = Inches(1 + i * 3.8)
        box = slide.shapes.add_shape(1, left, Inches(2), Inches(3.5), Inches(4))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = BLUE
        box.line.width = Pt(2)
        
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = BLUE
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(20)
        p2.font.color.rgb = GRAY
        p2.space_before = Pt(20)

    # --- 3. THE SOLUTION ---
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    txbox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.3), Inches(1))
    tf = txbox.text_frame
    p = tf.add_paragraph()
    p.text = "The Solution: Variant Sage"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE

    txbox2 = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.3), Inches(4))
    tf2 = txbox2.text_frame
    tf2.word_wrap = True
    
    p = tf2.add_paragraph()
    p.text = "An Agentic AI workflow that automates the heavy lifting of variant interpretation."
    p.font.size = Pt(32)
    p.font.color.rgb = TEAL
    
    bullets = [
        "Automated ACMG Classification: AI agents debate and classify variants (Pathogenic to Benign).",
        "Deep Literature Search: Instantly cross-references PubMed, ClinVar, and OMIM.",
        "Human-in-the-Loop: AI drafts the report, Genetic Counselors review and approve.",
        "Phenotype Mapping: Automatically matches patient HPO terms to variant databases."
    ]
    for b in bullets:
        p = tf2.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(24)
        p.font.color.rgb = WHITE
        p.space_before = Pt(15)

    # --- 4. THE MENA STRATEGY ---
    slide = prs.slides.add_slide(blank_layout)
    
    txbox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.3), Inches(1))
    tf = txbox.text_frame
    p = tf.add_paragraph()
    p.text = "Go-to-Market: The MENA Bundle Strategy"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Software + Hardware graphic
    left_box = slide.shapes.add_shape(1, Inches(1.5), Inches(2.5), Inches(4), Inches(3))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = TEAL
    left_box.text_frame.text = "Sequencing Hardware\n(Partner/White-label)"
    left_box.text_frame.paragraphs[0].font.size = Pt(28)
    left_box.text_frame.paragraphs[0].font.bold = True

    plus = slide.shapes.add_textbox(Inches(6), Inches(3.5), Inches(1), Inches(1))
    plus.text_frame.text = "+"
    plus.text_frame.paragraphs[0].font.size = Pt(60)
    plus.text_frame.paragraphs[0].font.bold = True
    plus.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    right_box = slide.shapes.add_shape(1, Inches(7.5), Inches(2.5), Inches(4), Inches(3))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = BLUE
    right_box.text_frame.text = "Variant Sage Software\n(AI Brain)"
    right_box.text_frame.paragraphs[0].font.size = Pt(28)
    right_box.text_frame.paragraphs[0].font.bold = True

    footer = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(11.3), Inches(1))
    f_p = footer.text_frame.add_paragraph()
    f_p.text = "Pitching a turnkey, integrated solution to hospitals rather than standalone software."
    f_p.font.size = Pt(24)
    f_p.font.color.rgb = GRAY
    f_p.alignment = PP_ALIGN.CENTER

    # --- 5. THE WORKFLOW ---
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BG
    bg.line.fill.background()

    txbox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.3), Inches(1))
    p = txbox.text_frame.add_paragraph()
    p.text = "End-to-End Hospital Workflow"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = NAVY

    steps = [
        ("1. Doctor", "Submits patient clinical history & HPO terms", Inches(1)),
        ("2. Lab", "Uploads VCF from sequencer", Inches(4)),
        ("3. AI Pipeline", "Variant Sage parses, filters, and classifies", Inches(7)),
        ("4. Counselor", "Reviews AI evidence & approves final report", Inches(10))
    ]

    for i, (title, desc, left) in enumerate(steps):
        # Circle
        circ = slide.shapes.add_shape(9, left + Inches(0.5), Inches(2.5), Inches(1.5), Inches(1.5)) # 9 is oval
        circ.fill.solid()
        circ.fill.fore_color.rgb = BLUE
        circ.line.fill.background()
        circ.text_frame.text = str(i+1)
        circ.text_frame.paragraphs[0].font.size = Pt(40)
        circ.text_frame.paragraphs[0].font.bold = True
        
        # Arrow (except last)
        if i < 3:
            arrow = slide.shapes.add_shape(33, left + Inches(2.2), Inches(3), Inches(0.6), Inches(0.5)) # right arrow
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GRAY
            arrow.line.fill.background()

        # Text
        tb = slide.shapes.add_textbox(left - Inches(0.25), Inches(4.5), Inches(3), Inches(2))
        tf = tb.text_frame
        tf.word_wrap = True
        p1 = tf.add_paragraph()
        p1.text = title.split(". ")[1]
        p1.font.size = Pt(24)
        p1.font.bold = True
        p1.font.color.rgb = NAVY
        p1.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(16)
        p2.font.color.rgb = GRAY
        p2.alignment = PP_ALIGN.CENTER

    # --- 6. METRICS SLIDE ---
    slide = prs.slides.add_slide(blank_layout)
    txbox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.3), Inches(1))
    p = txbox.text_frame.add_paragraph()
    p.text = "Why Variant Sage Wins"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = NAVY

    metrics = [
        ("94%", "Zero-Shot AI Success Rate", TEAL),
        ("10x", "Faster Turnaround Time", BLUE),
        ("-40%", "Reduction in VUS variants", NAVY)
    ]
    
    for i, (metric, desc, color) in enumerate(metrics):
        box = slide.shapes.add_shape(1, Inches(1 + i * 3.8), Inches(2.5), Inches(3.5), Inches(3.5))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = metric
        p.font.size = Pt(80)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(20)
        p2.alignment = PP_ALIGN.CENTER

    # Save
    prs.save('Variant_Sage_Pitch_Deck.pptx')
    print("Pitch deck created successfully.")

if __name__ == '__main__':
    create_deck()
